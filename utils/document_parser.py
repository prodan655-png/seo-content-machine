"""
Utility functions for extracting text from various document formats.
"""
from pathlib import Path
import io

def extract_text_from_pdf(file_content):
    """Extract text from PDF file."""
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except ImportError:
        return "[PDF parsing requires PyPDF2 - install with: pip install PyPDF2]"
    except Exception as e:
        return f"[Error extracting PDF: {e}]"

def extract_text_from_docx(file_content):
    """Extract text from DOCX file."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except ImportError:
        return "[DOCX parsing requires python-docx - install with: pip install python-docx]"
    except Exception as e:
        return f"[Error extracting DOCX: {e}]"

def extract_text_from_txt(file_content):
    """Extract text from TXT/MD file."""
    try:
        return file_content.decode('utf-8')
    except Exception as e:
        try:
            return file_content.decode('cp1251')  # Try Windows encoding
        except:
            return f"[Error extracting text: {e}]"

def extract_text_from_pptx(file_content):
    """Extract text from PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except ImportError:
        return "[PPTX parsing requires python-pptx - install with: pip install python-pptx]"
    except Exception as e:
        return f"[Error extracting PPTX: {e}]"

def extract_text_from_document(file_content, file_type):
    """
    Extract text from uploaded document based on file type.
    
    Args:
        file_content: Binary content of the file
        file_type: MIME type of the file
    
    Returns:
        Extracted text as string
    """
    if 'pdf' in file_type.lower():
        return extract_text_from_pdf(file_content)
    elif 'word' in file_type.lower() or 'docx' in file_type.lower():
        return extract_text_from_docx(file_content)
    elif 'presentation' in file_type.lower() or 'powerpoint' in file_type.lower() or 'pptx' in file_type.lower():
        return extract_text_from_pptx(file_content)
    elif 'text' in file_type.lower() or 'markdown' in file_type.lower():
        return extract_text_from_txt(file_content)
    else:
        return f"[Unsupported file type: {file_type}]"
